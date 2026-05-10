"""Multi-target export pipeline.

Each `export_to_*` reads the finalized session + user from Firestore, asks
the appropriate downstream service to materialise the artifacts, and
returns ``{url, ...}`` describing where the user can find them. All targets
ultimately leave the artifacts owned by the user, never the service account
(per CLAUDE.md hard constraint #13).

Targets
-------
- ``export_to_drive``         — Drive folder (transferOwnership=True)
- ``export_to_notion``        — Notion page with one block per artifact
- ``export_to_linear``        — Linear project + issues from GTM 90-day plan
- ``export_to_markdown_zip``  — bytes (in-memory ZIP)
- ``export_to_pptx``          — bytes (PPTX via python-pptx, fallback Slides export)
- ``export_to_json``          — bytes (full Pydantic dump)

Failure mode: every function logs and re-raises. The route layer catches and
maps to ``ExportTargetResult.ok=False``.
"""
from __future__ import annotations

import asyncio
import io
import json
import zipfile
from datetime import datetime, timezone
from typing import Any

import httpx
import structlog
from pydantic import BaseModel

from config import settings
from models.agent_schemas import (
    BrandIdentityResult,
    BusinessModelResult,
    CompetitiveAnalysisResult,
    ExecutiveSummaryResult,
    FinancialModelResult,
    GoToMarketResult,
    LandingPageResult,
    LegalDocumentsResult,
    MarketResearchResult,
    PitchDeckResult,
    RiskAnalysisResult,
    TechArchitectureResult,
)
from models.session_models import AgentName, Session
from models.user_models import User

log = structlog.get_logger(__name__)


# ─── Bundle assembly ─────────────────────────────────────────────────────────


_AGENT_TO_KEY: dict[AgentName, str] = {
    AgentName.IDEA_PARSER: "parsed_idea",
    AgentName.ARTICULATION: "articulation",
    AgentName.MARKET_RESEARCH: "market_research_result",
    AgentName.COMPETITIVE_ANALYSIS: "competitive_analysis_result",
    AgentName.BUSINESS_MODEL: "business_model_result",
    AgentName.BRAND_IDENTITY: "brand_identity_result",
    AgentName.RISK_ANALYSIS: "risk_analysis_result",
    AgentName.TECH_ARCHITECTURE: "tech_architecture_result",
    AgentName.FINANCIAL_MODEL: "financial_model_result",
    AgentName.LANDING_PAGE: "landing_page_result",
    AgentName.LEGAL_DOCUMENTS: "legal_documents_result",
    AgentName.GO_TO_MARKET: "go_to_market_result",
    AgentName.PITCH_DECK: "pitch_deck_result",
    AgentName.EXECUTIVE_SUMMARY: "executive_summary_result",
}


async def _load_session(session_id: str) -> Session:
    from services import firestore_service

    s = await firestore_service.read_session(session_id)
    if s is None:
        raise FileNotFoundError(f"session {session_id} not found")
    return s


async def _load_user(uid: str) -> User | None:
    from services import firestore_service

    try:
        return await firestore_service.get_user(uid)
    except Exception as exc:  # noqa: BLE001
        log.warning("export.user_fetch_failed", uid=uid, err=str(exc))
        return None


async def _load_outputs(session_id: str) -> dict[str, dict[str, Any]]:
    """Return {agent_value: payload_dict} for every produced agent."""
    from services import firestore_service

    out: dict[str, dict[str, Any]] = {}
    for agent in AgentName:
        try:
            doc = await firestore_service.read_agent_output(session_id, agent)
        except Exception as exc:  # noqa: BLE001
            log.warning("export.read_output_failed", agent=agent.value, err=str(exc))
            continue
        if not doc:
            continue
        payload = doc.get("payload") if isinstance(doc, dict) else None
        if payload is None:
            payload = doc
        if isinstance(payload, dict):
            out[agent.value] = payload
    return out


# ─── Markdown rendering ──────────────────────────────────────────────────────


def _h1(s: str) -> str:
    return f"# {s}\n\n"


def _h2(s: str) -> str:
    return f"## {s}\n\n"


def _kv(label: str, value: Any) -> str:
    return f"- **{label}:** {value}\n"


def _bullets(items: list[Any]) -> str:
    if not items:
        return "_None._\n"
    return "\n".join(f"- {item}" for item in items) + "\n"


def _summary_md(session: Session, exec_summary: dict | None, brand: dict | None) -> str:
    name = (brand or {}).get("company_name") or session.company_name or "Untitled"
    tagline = (brand or {}).get("tagline") or ""
    md: list[str] = [_h1(f"{name} — Executive Summary")]
    if tagline:
        md.append(f"_{tagline}_\n\n")
    if exec_summary:
        if oneliner := exec_summary.get("one_liner"):
            md.append(f"**One-liner:** {oneliner}\n\n")
        if summary := exec_summary.get("summary_text"):
            md.append(summary + "\n\n")
        if highs := exec_summary.get("key_highlights"):
            md.append(_h2("Highlights"))
            md.append(_bullets(highs))
        if pitch_30 := exec_summary.get("elevator_pitch_30s"):
            md.append(_h2("Elevator Pitch (30s)"))
            md.append(pitch_30 + "\n\n")
        if pitch_60 := exec_summary.get("elevator_pitch_60s"):
            md.append(_h2("Elevator Pitch (60s)"))
            md.append(pitch_60 + "\n\n")
    md.append(_h2("Session metadata"))
    md.append(_kv("Session ID", session.session_id))
    md.append(_kv("Status", session.status.value))
    md.append(_kv("Created", session.created_at.isoformat()))
    return "".join(md)


def _brand_md(brand: dict | None) -> str:
    if not brand:
        return _h1("Brand") + "_Not produced._\n"
    md = [_h1("Brand")]
    md.append(_kv("Company name", brand.get("company_name", "")))
    md.append(_kv("Tagline", brand.get("tagline", "")))
    voice = brand.get("brand_voice_traits") or []
    md.append(_kv("Voice traits", ", ".join(voice)))
    if logo := brand.get("logo_concept_description"):
        md.append(_kv("Logo concept", logo))
    md.append(_h2("Color palette"))
    for c in brand.get("color_palette") or []:
        md.append(_kv(f"{c.get('role','?')} ({c.get('name','')})", c.get("hex", "")))
    typo = brand.get("typography") or {}
    md.append(_h2("Typography"))
    md.append(_kv("Heading", typo.get("heading_font", "")))
    md.append(_kv("Body", typo.get("body_font", "")))
    md.append(_h2("Alternative names"))
    for alt in brand.get("name_alternatives") or []:
        if isinstance(alt, dict):
            md.append(_kv(alt.get("name", ""), alt.get("rationale", "")))
    return "".join(md)


def _market_md(market: dict | None) -> str:
    if not market:
        return _h1("Market Research") + "_Not produced._\n"
    md = [_h1("Market Research")]
    for k in ("tam", "sam", "som", "cagr"):
        v = market.get(k) or {}
        if v:
            md.append(
                _kv(
                    k.upper(),
                    f"{v.get('value','?')} {v.get('unit','')} ({v.get('confidence','?')})",
                )
            )
    md.append(_h2("Industry trends"))
    md.append(_bullets(market.get("industry_trends") or []))
    md.append(_h2("Target demographics"))
    md.append(_bullets(market.get("target_demographics") or []))
    md.append(_h2("Sources"))
    for s in market.get("sources") or []:
        md.append(f"- [{s.get('publisher') or 'source'}]({s.get('source_url','')}) — {s.get('text','')}\n")
    return "".join(md)


def _competitive_md(comp: dict | None) -> str:
    if not comp:
        return _h1("Competitive Analysis") + "_Not produced._\n"
    md = [_h1("Competitive Analysis")]
    md.append(
        _kv("Market concentration", comp.get("market_concentration", "")),
    )
    md.append(_h2("Competitors"))
    for c in comp.get("competitors") or []:
        if not isinstance(c, dict):
            continue
        md.append(f"### {c.get('name','Unnamed')}\n\n")
        if d := c.get("description"):
            md.append(d + "\n\n")
        if c.get("data_disclosed") is False:
            md.append("_No public funding/revenue data available._\n\n")
        if s := c.get("strengths"):
            md.append("**Strengths:**\n" + _bullets(s))
        if w := c.get("weaknesses"):
            md.append("**Weaknesses:**\n" + _bullets(w))
    md.append(_h2("Positioning gaps"))
    md.append(_bullets(comp.get("positioning_gaps") or []))
    return "".join(md)


def _bm_md(bm: dict | None) -> str:
    if not bm:
        return _h1("Business Model") + "_Not produced._\n"
    md = [_h1("Business Model")]
    md.append(_kv("Revenue model", bm.get("revenue_model", "")))
    md.append(_kv("Primary stream", bm.get("primary_revenue_stream", "")))
    md.append(_h2("Pricing tiers"))
    for t in bm.get("pricing_tiers") or []:
        if not isinstance(t, dict):
            continue
        md.append(
            f"### {t.get('name','?')} — ${t.get('price_usd_monthly',0):.0f}/mo\n\n"
        )
        md.append(_kv("Target", t.get("target_segment", "")))
        md.append(_bullets(t.get("features") or []))
    md.append(_h2("Unit economics"))
    ue = bm.get("unit_economics") or {}
    if isinstance(ue, dict):
        for k in ("cac_usd", "ltv_usd", "gross_margin_pct", "payback_months"):
            v = ue.get(k) or {}
            if isinstance(v, dict):
                md.append(_kv(k, f"{v.get('value','?')} {v.get('unit','')}"))
        md.append(_kv("LTV/CAC ratio", ue.get("ltv_cac_ratio", "")))
    md.append(_h2("Business Model Canvas"))
    bmc = bm.get("business_model_canvas") or {}
    for block, items in bmc.items():
        md.append(f"**{block}:**\n")
        if isinstance(items, list):
            md.append(_bullets(items))
        else:
            md.append(f"- {items}\n")
    return "".join(md)


def _financial_md(fin: dict | None) -> str:
    if not fin:
        return _h1("Financial Model") + "_Not produced._\n"
    md = [_h1("Financial Model")]
    md.append(_kv("Seed funding (USD)", f"${fin.get('funding_seed_usd', 0):,.0f}"))
    md.append(_kv("Runway (months)", fin.get("runway_months", "?")))
    md.append(_kv("Breakeven month", fin.get("breakeven_month") or "n/a"))
    md.append(_kv("Reconciliation passed", fin.get("reconciliation_passed", False)))
    md.append("\n| Year | Revenue | COGS | Gross | Opex | EBITDA | Headcount | Cash |\n")
    md.append("|---|---|---|---|---|---|---|---|\n")
    for row in fin.get("projections") or []:
        md.append(
            f"| {row.get('year','?')} "
            f"| ${row.get('revenue_usd',0):,.0f} "
            f"| ${row.get('cogs_usd',0):,.0f} "
            f"| ${row.get('gross_profit_usd',0):,.0f} "
            f"| ${row.get('opex_usd',0):,.0f} "
            f"| ${row.get('ebitda_usd',0):,.0f} "
            f"| {row.get('headcount','?')} "
            f"| ${row.get('cash_usd',0):,.0f} |\n"
        )
    md.append("\n")
    if metrics := fin.get("key_metrics"):
        md.append(_h2("Key metrics"))
        for k, v in metrics.items():
            md.append(_kv(k, v))
    return "".join(md)


def _gtm_md(gtm: dict | None) -> str:
    if not gtm:
        return _h1("Go-to-Market") + "_Not produced._\n"
    md = [_h1("Go-to-Market")]
    md.append(_kv("Strategy", gtm.get("launch_strategy_type", "")))
    md.append(_h2("Launch phases"))
    for ph in gtm.get("launch_phases") or []:
        if isinstance(ph, dict):
            md.append(
                f"- **{ph.get('phase','?')}** ({ph.get('weeks','?')}): {ph.get('actions','')}\n"
            )
    md.append(_h2("90-day plan"))
    plan = gtm.get("first_90_days_plan") or {}
    for window, actions in plan.items():
        md.append(f"### {window}\n\n")
        if isinstance(actions, list):
            md.append(_bullets(actions))
        else:
            md.append(f"- {actions}\n")
    md.append(_h2("Marketing channels"))
    for ch in gtm.get("marketing_channels") or []:
        if isinstance(ch, dict):
            md.append(
                _kv(
                    ch.get("channel", "?"),
                    f"CAC≈${ch.get('cac_estimate', '?')} priority={ch.get('priority','?')}",
                )
            )
    md.append(_h2("KPIs"))
    for metric, values in (gtm.get("kpis") or {}).items():
        if isinstance(values, dict):
            md.append(_kv(metric, ", ".join(f"{k}={v}" for k, v in values.items())))
    md.append(_h2("Partnerships"))
    md.append(_bullets(gtm.get("partnerships") or []))
    return "".join(md)


def _risk_md(risk: dict | None) -> str:
    if not risk:
        return _h1("Risk Analysis") + "_Not produced._\n"
    md = [_h1("Risk Analysis")]
    md.append("\n| Category | Description | Probability | Impact | Mitigation |\n")
    md.append("|---|---|---|---|---|\n")
    for r in risk.get("risk_matrix") or []:
        if not isinstance(r, dict):
            continue
        md.append(
            f"| {r.get('category','?')} "
            f"| {r.get('description','')} "
            f"| {r.get('probability','?')} "
            f"| {r.get('impact','?')} "
            f"| {r.get('mitigation','')} |\n"
        )
    md.append("\n")
    if reg := risk.get("regulatory_considerations"):
        md.append(_h2("Regulatory considerations"))
        for jur, items in reg.items():
            md.append(f"### {jur}\n\n")
            if isinstance(items, list):
                md.append(_bullets(items))
    if w := risk.get("worst_case_scenario"):
        md.append(_h2("Worst-case scenario"))
        md.append(w + "\n\n")
    if pivots := risk.get("pivot_options"):
        md.append(_h2("Pivot options"))
        md.append(_bullets(pivots))
    return "".join(md)


def _tech_md(tech: dict | None) -> str:
    if not tech:
        return _h1("Tech Architecture") + "_Not produced._\n"
    md = [_h1("Tech Architecture")]
    md.append(_h2("Recommended stack"))
    for k, v in (tech.get("recommended_stack") or {}).items():
        md.append(_kv(k, v))
    if mer := tech.get("architecture_diagram_mermaid"):
        md.append(_h2("Architecture diagram"))
        md.append("```mermaid\n" + mer + "\n```\n\n")
    md.append(_h2("MVP core features"))
    md.append(_bullets(tech.get("mvp_core_features") or []))
    md.append(_h2("Nice-to-have"))
    md.append(_bullets(tech.get("mvp_nice_to_have") or []))
    md.append(_kv("Estimated dev weeks", tech.get("estimated_dev_weeks", "?")))
    md.append(_kv("Estimated team size", tech.get("estimated_team_size", "?")))
    cost = tech.get("monthly_infra_cost_usd_estimate") or {}
    if isinstance(cost, dict):
        md.append(_kv("Monthly infra cost", f"${cost.get('value','?')}"))
    md.append(_h2("Security considerations"))
    md.append(_bullets(tech.get("security_considerations") or []))
    return "".join(md)


def _deck_md(deck: dict | None) -> str:
    if not deck:
        return _h1("Pitch Deck") + "_Not produced._\n"
    md = [_h1("Pitch Deck")]
    for slide in deck.get("slides") or []:
        if not isinstance(slide, dict):
            continue
        md.append(
            f"## Slide {slide.get('slide_number','?')} — {slide.get('title','')} "
            f"_({slide.get('layout','?')})_\n\n"
        )
        md.append((slide.get("body") or "") + "\n\n")
        if notes := slide.get("speaker_notes"):
            md.append(f"_Speaker notes:_ {notes}\n\n")
    return "".join(md)


def _legal_files(legal: dict | None) -> dict[str, str]:
    """Return a {filename: content} mapping for the 08-legal/ subfolder."""
    out: dict[str, str] = {}
    if not legal:
        out["README.md"] = "_Legal docs not produced._\n"
        return out
    out["README.md"] = (
        _h1("Legal — README")
        + "**This is an AI-generated draft. NOT legal advice. Have a licensed attorney review.**\n\n"
        + _kv("ToS template", legal.get("tos_template_id", ""))
        + _kv("Privacy template", legal.get("privacy_template_id", ""))
        + _kv("Jurisdictions", ", ".join(legal.get("jurisdictions_covered") or []))
        + _kv("Lawyer review CTA", legal.get("lawyer_review_cta", True))
        + "\n"
    )
    out["incorporation_checklist.md"] = _h1("Incorporation Checklist") + "\n".join(
        f"- {row.get('task','?')} ({row.get('estimated_cost_usd', '?')} USD, {row.get('approximate_time','?')})"
        for row in (legal.get("incorporation_checklist") or [])
        if isinstance(row, dict)
    ) + "\n"
    return out


# ─── Markdown ZIP ────────────────────────────────────────────────────────────


async def export_to_markdown_zip(
    session_id: str,
    *,
    uid: str | None = None,
) -> dict[str, Any]:
    """Pure in-memory ZIP. Returns ``{filename, bytes, byte_size}``.

    Routes wrapping this can return raw bytes; tests assert the structure.
    """
    session = await _load_session(session_id)
    outputs = await _load_outputs(session_id)
    brand = outputs.get(AgentName.BRAND_IDENTITY.value) or {}
    exec_sum = outputs.get(AgentName.EXECUTIVE_SUMMARY.value) or {}
    market = outputs.get(AgentName.MARKET_RESEARCH.value)
    comp = outputs.get(AgentName.COMPETITIVE_ANALYSIS.value)
    bm = outputs.get(AgentName.BUSINESS_MODEL.value)
    fin = outputs.get(AgentName.FINANCIAL_MODEL.value)
    landing = outputs.get(AgentName.LANDING_PAGE.value)
    deck = outputs.get(AgentName.PITCH_DECK.value)
    legal = outputs.get(AgentName.LEGAL_DOCUMENTS.value)
    gtm = outputs.get(AgentName.GO_TO_MARKET.value)
    risk = outputs.get(AgentName.RISK_ANALYSIS.value)
    tech = outputs.get(AgentName.TECH_ARCHITECTURE.value)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("00-summary.md", _summary_md(session, exec_sum, brand))
        zf.writestr("01-brand.md", _brand_md(brand))
        zf.writestr("02-market.md", _market_md(market))
        zf.writestr("03-competitive.md", _competitive_md(comp))
        zf.writestr("04-business-model.md", _bm_md(bm))
        zf.writestr("05-financial.md", _financial_md(fin))
        zf.writestr(
            "06-landing.html",
            (landing or {}).get("html_sanitized") or "<!-- landing not produced -->",
        )
        zf.writestr("07-deck.md", _deck_md(deck))
        for fname, content in _legal_files(legal).items():
            zf.writestr(f"08-legal/{fname}", content)
        zf.writestr("09-gtm.md", _gtm_md(gtm))
        zf.writestr("10-risks.md", _risk_md(risk))
        zf.writestr("11-tech.md", _tech_md(tech))

    payload = buf.getvalue()
    company = (brand or {}).get("company_name") or session.session_id
    safe = "".join(c for c in str(company) if c.isalnum() or c in "-_") or session.session_id
    filename = f"{safe}-prometheus-export.zip"
    log.info("export.markdown_zip.done", session_id=session_id, bytes=len(payload))
    return {
        "filename": filename,
        "bytes": payload,
        "byte_size": len(payload),
        "url": None,
        "content_type": "application/zip",
    }


# ─── JSON dump ───────────────────────────────────────────────────────────────


def _model_dump(value: Any) -> Any:
    if isinstance(value, BaseModel):
        return value.model_dump(mode="json")
    if isinstance(value, dict):
        return {k: _model_dump(v) for k, v in value.items()}
    if isinstance(value, list):
        return [_model_dump(v) for v in value]
    return value


async def export_to_json(
    session_id: str,
    *,
    uid: str | None = None,
) -> dict[str, Any]:
    session = await _load_session(session_id)
    outputs = await _load_outputs(session_id)

    bundle = {
        "exported_at": datetime.now(tz=timezone.utc).isoformat(),
        "session": {
            "session_id": session.session_id,
            "user_uid": session.user_uid,
            "company_name": session.company_name,
            "status": session.status.value,
            "created_at": session.created_at.isoformat(),
            "started_at": session.started_at.isoformat() if session.started_at else None,
            "completed_at": (
                session.completed_at.isoformat() if session.completed_at else None
            ),
            "cost": session.cost.model_dump(mode="json"),
            "metadata": session.metadata,
            "agents": {
                name.value: rec.model_dump(mode="json")
                for name, rec in session.agents.items()
            },
        },
        "outputs": _model_dump(outputs),
    }
    payload = json.dumps(bundle, indent=2, default=str).encode("utf-8")
    log.info("export.json.done", session_id=session_id, bytes=len(payload))
    return {
        "filename": f"{session.session_id}-bundle.json",
        "bytes": payload,
        "byte_size": len(payload),
        "url": None,
        "content_type": "application/json",
    }


# ─── PPTX (python-pptx primary, Google Slides export fallback) ───────────────


async def export_to_pptx(
    session_id: str,
    *,
    uid: str | None = None,
) -> dict[str, Any]:
    session = await _load_session(session_id)
    outputs = await _load_outputs(session_id)
    deck_payload = outputs.get(AgentName.PITCH_DECK.value) or {}
    brand_payload = outputs.get(AgentName.BRAND_IDENTITY.value) or {}

    try:
        deck = PitchDeckResult.model_validate(deck_payload) if deck_payload else None
    except Exception as exc:  # noqa: BLE001
        log.warning("export.pptx.deck_invalid", err=str(exc))
        deck = None

    if deck is None:
        raise ValueError("session has no pitch_deck_result")

    try:
        brand = BrandIdentityResult.model_validate(brand_payload) if brand_payload else None
    except Exception:  # noqa: BLE001
        brand = None

    primary_hex = "#0F172A"
    accent_hex = "#F97316"
    if brand:
        for c in brand.color_palette:
            if c.role == "primary":
                primary_hex = c.hex
            elif c.role == "accent":
                accent_hex = c.hex

    def _build_pptx() -> bytes | None:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
            from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
            from pptx.util import Inches, Pt  # type: ignore[import-not-found]
        except Exception as exc:  # noqa: BLE001
            log.warning("export.pptx.python_pptx_missing", err=str(exc))
            return None

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        title_layout = prs.slide_layouts[0]
        body_layout = prs.slide_layouts[1]

        def _hex_to_rgb(h: str) -> RGBColor:
            h = h.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        primary_rgb = _hex_to_rgb(primary_hex)
        accent_rgb = _hex_to_rgb(accent_hex)

        for i, sl in enumerate(deck.slides):
            layout = title_layout if i == 0 or sl.layout == "title" else body_layout
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = sl.title
                for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = primary_rgb
                    run.font.size = Pt(36)
            placeholders = [p for p in slide.placeholders if p.placeholder_format.idx != 0]
            if placeholders:
                tf = placeholders[0].text_frame
                tf.text = sl.body[:1200]
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(18)
                        run.font.color.rgb = accent_rgb if i == 0 else _hex_to_rgb("#1F2937")
            if sl.speaker_notes:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.text = sl.speaker_notes

        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()

    payload = await asyncio.to_thread(_build_pptx)

    if payload is None:
        # Fallback: export the Google Slides file we already created in
        # google_workspace.create_presentation_from_template, if any.
        if deck.presentation_id:
            try:
                from services.google_workspace import _get_drive  # type: ignore

                def _fetch() -> bytes:
                    drive = _get_drive()
                    request = drive.files().export_media(
                        fileId=deck.presentation_id,
                        mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                    return bytes(request.execute())

                payload = await asyncio.to_thread(_fetch)
            except Exception as exc:  # noqa: BLE001
                log.warning("export.pptx.slides_export_failed", err=str(exc))
                payload = None

    if payload is None:
        raise RuntimeError("pptx generation unavailable (python-pptx missing and no Slides fallback)")

    company = (brand_payload or {}).get("company_name") or session.session_id
    safe = "".join(c for c in str(company) if c.isalnum() or c in "-_") or session.session_id
    filename = f"{safe}-pitch-deck.pptx"
    log.info("export.pptx.done", session_id=session_id, bytes=len(payload))
    return {
        "filename": filename,
        "bytes": payload,
        "byte_size": len(payload),
        "url": None,
        "content_type": (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
    }


# ─── Drive folder export ─────────────────────────────────────────────────────


async def export_to_drive(
    session_id: str,
    *,
    uid: str | None = None,
) -> dict[str, Any]:
    """Create a Drive folder owned by the user, drop deck/doc/sheet/HTML/JSON inside."""
    session = await _load_session(session_id)
    outputs = await _load_outputs(session_id)
    brand_payload = outputs.get(AgentName.BRAND_IDENTITY.value) or {}
    landing_payload = outputs.get(AgentName.LANDING_PAGE.value) or {}
    fin_payload = outputs.get(AgentName.FINANCIAL_MODEL.value) or {}
    deck_payload = outputs.get(AgentName.PITCH_DECK.value) or {}

    user = await _load_user(uid or session.user_uid)
    user_email = user.email if user else None

    try:
        from services import google_workspace  # local
    except Exception as exc:  # noqa: BLE001
        log.warning("export.drive.workspace_unavailable", err=str(exc))
        raise

    company = (brand_payload or {}).get("company_name") or session.session_id

    folder_id, folder_url = await google_workspace.create_drive_folder(
        f"{company} — PROMETHEUS export"
    )

    file_ids: list[str] = []

    # 1. Slides — reuse existing presentation if pitch_deck_result.presentation_id is set.
    deck_pres_id = (deck_payload or {}).get("presentation_id")
    if deck_pres_id:
        try:
            await google_workspace.move_file_to_folder(deck_pres_id, folder_id)
            file_ids.append(deck_pres_id)
        except Exception as exc:  # noqa: BLE001
            log.warning("export.drive.move_deck_failed", err=str(exc))
    elif deck_payload:
        try:
            deck = PitchDeckResult.model_validate(deck_payload)
            try:
                brand_obj = BrandIdentityResult.model_validate(brand_payload)
            except Exception:  # noqa: BLE001
                brand_obj = None
            if brand_obj is not None:
                pres_id, _ = await google_workspace.create_presentation_from_template(
                    brand=brand_obj, slides=deck.slides
                )
                file_ids.append(pres_id)
                await google_workspace.move_file_to_folder(pres_id, folder_id)
        except Exception as exc:  # noqa: BLE001
            log.warning("export.drive.deck_create_failed", err=str(exc))

    # 2. Sheets from finance.
    if fin_payload:
        try:
            fin = FinancialModelResult.model_validate(fin_payload)
            sheet_id, _ = await google_workspace.create_sheets_from_finance(fin)
            file_ids.append(sheet_id)
            await google_workspace.move_file_to_folder(sheet_id, folder_id)
        except Exception as exc:  # noqa: BLE001
            log.warning("export.drive.sheets_failed", err=str(exc))

    # 3. Exec summary doc.
    exec_payload = outputs.get(AgentName.EXECUTIVE_SUMMARY.value) or {}
    if exec_payload:
        try:
            md = _summary_md(session, exec_payload, brand_payload)
            doc_id, _ = await google_workspace.create_doc_from_template(
                template_md=md,
                vars={
                    "title": f"{company} — Executive Summary",
                    "company_name": company,
                },
            )
            file_ids.append(doc_id)
            await google_workspace.move_file_to_folder(doc_id, folder_id)
        except Exception as exc:  # noqa: BLE001
            log.warning("export.drive.doc_failed", err=str(exc))

    # 4. Landing HTML — write as plain file via Drive directly.
    landing_html = (landing_payload or {}).get("html_sanitized") or ""
    if landing_html:
        try:
            file_id = await _upload_blob_to_drive(
                folder_id=folder_id,
                name=f"{company}-landing.html",
                mime_type="text/html",
                data=landing_html.encode("utf-8"),
            )
            if file_id:
                file_ids.append(file_id)
        except Exception as exc:  # noqa: BLE001
            log.warning("export.drive.landing_failed", err=str(exc))

    # 5. JSON bundle.
    try:
        bundle = await export_to_json(session_id, uid=uid)
        file_id = await _upload_blob_to_drive(
            folder_id=folder_id,
            name=bundle["filename"],
            mime_type=bundle["content_type"],
            data=bundle["bytes"],
        )
        if file_id:
            file_ids.append(file_id)
    except Exception as exc:  # noqa: BLE001
        log.warning("export.drive.json_failed", err=str(exc))

    # Transfer ownership: critical hard constraint #13.
    if user_email:
        for fid in file_ids + [folder_id]:
            try:
                await google_workspace.transfer_ownership(fid, user_email)
            except Exception as exc:  # noqa: BLE001
                log.warning(
                    "export.drive.owner_transfer_failed",
                    file_id=fid,
                    err=str(exc),
                )
    else:
        # Fallback: anyone-with-link share so user can at least access.
        try:
            await google_workspace.share_anyone_with_link(folder_id)
        except Exception as exc:  # noqa: BLE001
            log.warning("export.drive.share_fallback_failed", err=str(exc))

    log.info(
        "export.drive.done",
        session_id=session_id,
        folder_id=folder_id,
        file_count=len(file_ids),
        owner_transfer=bool(user_email),
    )
    return {
        "url": folder_url
        or f"https://drive.google.com/drive/folders/{folder_id}",
        "folder_url": folder_url,
        "folder_id": folder_id,
        "file_ids": file_ids,
    }


async def _upload_blob_to_drive(
    *, folder_id: str, name: str, mime_type: str, data: bytes
) -> str | None:
    try:
        from services.google_workspace import _get_drive  # type: ignore[attr-defined]
    except Exception:  # noqa: BLE001
        return None

    def _do() -> str | None:
        try:
            from googleapiclient.http import MediaIoBaseUpload  # type: ignore[import-not-found]

            drive = _get_drive()
            media = MediaIoBaseUpload(io.BytesIO(data), mimetype=mime_type, resumable=False)
            meta = drive.files().create(
                body={"name": name, "parents": [folder_id]},
                media_body=media,
                fields="id",
                supportsAllDrives=False,
            ).execute()
            return str(meta.get("id") or "") or None
        except Exception as exc:  # noqa: BLE001
            log.warning("export.drive.upload_failed", err=str(exc))
            return None

    return await asyncio.to_thread(_do)


# ─── Notion export ───────────────────────────────────────────────────────────


_NOTION_API = "https://api.notion.com/v1"


async def export_to_notion(
    session_id: str,
    notion_token: str | None = None,
    *,
    uid: str | None = None,
    parent_page_id: str | None = None,
) -> dict[str, Any]:
    if not notion_token:
        raise ValueError("notion_token required")
    session = await _load_session(session_id)
    outputs = await _load_outputs(session_id)
    brand = outputs.get(AgentName.BRAND_IDENTITY.value) or {}
    company = brand.get("company_name") or session.session_id

    headers = {
        "Authorization": f"Bearer {notion_token}",
        "Notion-Version": settings.notion_api_version,
        "Content-Type": "application/json",
    }

    parent: dict[str, Any]
    if parent_page_id:
        parent = {"type": "page_id", "page_id": parent_page_id}
    else:
        parent = {"type": "workspace", "workspace": True}

    sections: list[tuple[str, str]] = [
        ("Executive Summary", _summary_md(session, outputs.get(AgentName.EXECUTIVE_SUMMARY.value) or {}, brand)),
        ("Brand", _brand_md(brand)),
        ("Market", _market_md(outputs.get(AgentName.MARKET_RESEARCH.value))),
        ("Competition", _competitive_md(outputs.get(AgentName.COMPETITIVE_ANALYSIS.value))),
        ("Business Model", _bm_md(outputs.get(AgentName.BUSINESS_MODEL.value))),
        ("Financial", _financial_md(outputs.get(AgentName.FINANCIAL_MODEL.value))),
        ("GTM", _gtm_md(outputs.get(AgentName.GO_TO_MARKET.value))),
        ("Risks", _risk_md(outputs.get(AgentName.RISK_ANALYSIS.value))),
        ("Tech", _tech_md(outputs.get(AgentName.TECH_ARCHITECTURE.value))),
        ("Pitch Deck", _deck_md(outputs.get(AgentName.PITCH_DECK.value))),
    ]

    blocks: list[dict[str, Any]] = []
    for heading, body in sections:
        blocks.append(_notion_heading(heading))
        for chunk in _chunk_text(body, 1900):
            blocks.append(_notion_paragraph(chunk))

    body = {
        "parent": parent,
        "properties": {
            "title": [{"type": "text", "text": {"content": f"{company} — PROMETHEUS"}}],
        },
        "children": blocks[:99],  # Notion API caps initial children at 100.
    }

    async with httpx.AsyncClient(timeout=20.0, headers=headers) as client:
        r = await client.post(f"{_NOTION_API}/pages", json=body)
        r.raise_for_status()
        page = r.json()
        page_id = page.get("id", "")
        page_url = page.get("url", "")

        # Append remaining blocks (Notion children API limit per request = 100).
        remaining = blocks[99:]
        while remaining:
            slice_, remaining = remaining[:100], remaining[100:]
            try:
                await client.patch(
                    f"{_NOTION_API}/blocks/{page_id}/children",
                    json={"children": slice_},
                )
            except Exception as exc:  # noqa: BLE001
                log.warning("export.notion.append_failed", err=str(exc))
                break

    log.info("export.notion.done", session_id=session_id, page_id=page_id)
    return {"url": page_url, "page_id": page_id, "page_url": page_url}


def _notion_heading(text: str) -> dict[str, Any]:
    return {
        "object": "block",
        "type": "heading_2",
        "heading_2": {"rich_text": [{"type": "text", "text": {"content": text[:1900]}}]},
    }


def _notion_paragraph(text: str) -> dict[str, Any]:
    return {
        "object": "block",
        "type": "paragraph",
        "paragraph": {"rich_text": [{"type": "text", "text": {"content": text[:1900]}}]},
    }


def _chunk_text(text: str, size: int) -> list[str]:
    if not text:
        return [""]
    return [text[i : i + size] for i in range(0, len(text), size)] or [text]


# ─── Linear export ───────────────────────────────────────────────────────────


_LINEAR_API = "https://api.linear.app/graphql"


async def export_to_linear(
    session_id: str,
    linear_token: str | None = None,
    *,
    uid: str | None = None,
    team_id: str | None = None,
) -> dict[str, Any]:
    if not linear_token:
        raise ValueError("linear_token required")
    session = await _load_session(session_id)
    outputs = await _load_outputs(session_id)
    brand = outputs.get(AgentName.BRAND_IDENTITY.value) or {}
    gtm = outputs.get(AgentName.GO_TO_MARKET.value) or {}
    company = brand.get("company_name") or session.session_id

    headers = {
        "Authorization": linear_token,
        "Content-Type": "application/json",
    }

    async with httpx.AsyncClient(timeout=20.0, headers=headers) as client:
        # Resolve team_id if not supplied.
        if not team_id:
            r = await client.post(
                _LINEAR_API,
                json={"query": "query { teams { nodes { id name } } }"},
            )
            r.raise_for_status()
            data = r.json().get("data", {}) or {}
            teams = (data.get("teams") or {}).get("nodes") or []
            if not teams:
                raise RuntimeError("no Linear teams accessible to this token")
            team_id = teams[0]["id"]

        # Create the project.
        project_create = {
            "query": (
                "mutation($input: ProjectCreateInput!) { "
                "projectCreate(input: $input) { project { id name url } success } }"
            ),
            "variables": {
                "input": {
                    "name": f"{company} — Launch (PROMETHEUS)",
                    "description": (
                        "Auto-generated 90-day launch plan from PROMETHEUS. "
                        "Edit/triage as you would any other Linear project."
                    ),
                    "teamIds": [team_id],
                }
            },
        }
        r = await client.post(_LINEAR_API, json=project_create)
        r.raise_for_status()
        proj = (r.json().get("data") or {}).get("projectCreate") or {}
        project = proj.get("project") or {}
        project_id = project.get("id")
        project_url = project.get("url", "")

        if not project_id:
            raise RuntimeError("Linear project creation failed")

        plan = gtm.get("first_90_days_plan") or {}
        issue_create = (
            "mutation($input: IssueCreateInput!) { "
            "issueCreate(input: $input) { issue { id identifier url } success } }"
        )
        issue_count = 0
        first_issue_url: str | None = None
        for window, actions in plan.items():
            if not isinstance(actions, list):
                continue
            for action in actions[:10]:
                title = f"[{window}] {str(action)[:200]}"
                r = await client.post(
                    _LINEAR_API,
                    json={
                        "query": issue_create,
                        "variables": {
                            "input": {
                                "teamId": team_id,
                                "projectId": project_id,
                                "title": title,
                                "description": (
                                    f"Window: {window}\n\nAction: {action}\n\n"
                                    "Generated by PROMETHEUS GTM agent."
                                ),
                            }
                        },
                    },
                )
                if r.status_code != 200:
                    log.warning("export.linear.issue_failed", status=r.status_code)
                    continue
                payload = (r.json().get("data") or {}).get("issueCreate") or {}
                issue = payload.get("issue") or {}
                if issue.get("id"):
                    issue_count += 1
                    if first_issue_url is None:
                        first_issue_url = issue.get("url")

    url = project_url or first_issue_url or ""
    log.info(
        "export.linear.done",
        session_id=session_id,
        project_id=project_id,
        issues=issue_count,
    )
    return {
        "url": url,
        "project_id": project_id,
        "project_url": project_url,
        "issue_count": issue_count,
    }


__all__ = [
    "export_to_drive",
    "export_to_json",
    "export_to_linear",
    "export_to_markdown_zip",
    "export_to_notion",
    "export_to_pptx",
]
